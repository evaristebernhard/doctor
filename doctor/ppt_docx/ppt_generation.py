'''将大模型生成的json数据转换为ppt，可修改代码自定义ppt的样式'''
# 导入标准库
import datetime  # 日期时间处理模块，虽然导入了但未使用
import hashlib  # 哈希算法模块，用于生成文件名
import os  # 操作系统接口模块，用于文件路径操作
import time  # 时间模块，用于生成时间戳


# 导入第三方库
from pptx.oxml.ns import qn  # PowerPoint XML命名空间模块，用于处理中文显示
from typing import Dict  # 类型提示，用于字典类型注解
from pptx import Presentation  # Python-PPTX库，用于创建和操作PPT文件

# 导入项目模块
from env import get_app_root  # 获取应用根目录的函数



# 定义PPT文件输出目录路径，将PPT文件保存在应用根目录下的data/cache/ppt文件夹中
_OUTPUT_DIR = os.path.join(get_app_root(), "data/cache/ppt")

# 如果文件夹路径不存在，先创建该目录
if not os.path.exists(_OUTPUT_DIR):
    # 递归创建目录，确保所有父目录都存在
    os.makedirs(_OUTPUT_DIR)

def get_file_path(text):
    """
    生成唯一的文件路径
    
    Args:
        text (str): 用于生成文件名的文本
        
    Returns:
        str: 完整的文件路径
    """
    # 使用SHA256哈希算法对文本进行哈希，生成唯一的文件名
    file_name = hashlib.sha256(text.encode("utf-8")).hexdigest()  ## 也可以使用uuid
    # 返回完整的文件路径，文件扩展名为.pptx
    return os.path.join(_OUTPUT_DIR, f"{file_name}.pptx")

def generate(ppt_content: Dict) -> str:
    """
    根据内容字典生成PPT文件
    
    Args:
        ppt_content (Dict): 包含PPT内容的字典，应包含标题和页面内容
        
    Returns:
        str: 生成的PPT文件路径
    """
    # 创建一个新的演示文稿对象
    ppt = Presentation()

    # PPT首页 - 使用第0号布局（标题和副标题布局）
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])  # title&subtitle layout
    # 设置首页标题，从ppt_content字典中获取title字段
    slide.placeholders[0].text = ppt_content["title"]
    # 设置首页副标题，固定为"--来自「赛博华佗」"
    slide.placeholders[1].text = "--来自「赛博华佗」"

    # 内容页 - 遍历所有页面内容
    print(f"总共{len(ppt_content['pages'])}页")
    # 遍历pages数组中的每个页面
    for i, page in enumerate(ppt_content["pages"]):
        # 打印当前页面信息到控制台
        print("生成第%d页:%s" % (i + 1, page["title"]))
        # 添加新幻灯片，使用第1号布局（标题和内容布局）
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # title&content layout
        
        # 标题 - 设置当前页面的标题
        slide.placeholders[0].text = page["title"]
        # 正文 - 获取幻灯片内容占位符的文本框架对象
        text_frame = slide.placeholders[1].text_frame  # 获取文本框的text_frame对象
        

        # 遍历当前页面的内容项
        for sub_content in page["content"]:
            # 打印子内容到控制台
            print(sub_content)
            
            # 一级正文 - 添加段落作为一级正文
            sub_title = text_frame.add_paragraph()
            # 设置段落文本和层级（level=2表示一级正文）
            sub_title.text, sub_title.level = sub_content["title"], 2
            
            # 二级正文 - 添加段落作为二级正文
            sub_description = text_frame.add_paragraph()
            # 设置段落文本和层级（level=3表示二级正文）
            sub_description.text, sub_description.level = sub_content["description"], 3
            
    # 生成输出文件路径，使用当前时间戳确保文件名唯一
    _output_file = get_file_path(str(time.time()))
    # 保存PPT文件到指定路径
    ppt.save(_output_file)

    # 返回生成的PPT文件路径
    return _output_file